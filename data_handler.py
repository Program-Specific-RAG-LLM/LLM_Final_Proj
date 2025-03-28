"""
Class that manages data for the RAG LLM. It prepares the data for the model, which should be populated in the data folder.

The data that will be accepeted in the following formats:
- PDF
- DOCX
- TXT
- PPTX
- NXML (specifically from statpearls)

Other formats will be ignored. Once cleaned the data will be extracted to a unified format (txt) the vectorized to be loaded into a vector DB.
"""

# Standard imports
import os
import re
import json
from typing import Tuple
from pathlib import Path

# Internal imports
from utils import embed_text
from const import PATH_TO_DATA, PATH_TO_CLEANED_DATA, PATH_TO_VECTORIZED_DATA

# External imports
import fitz
import docx2txt
import lxml.etree as et
from tqdm import tqdm
from pptx import Presentation


class DataHandler:
    """
    Class that handles the data for the RAG LLM.

    Attributes:
    - data_path: str, path to the data folder
    - data: list, list of paths to the data files
    - file_types: list, list of accepted file types

    Methods:
    - load_data: loads the data from the data folder
    - clean_data: cleans the data
    """

    def __init__(
        self,
        data_path: Path = Path(PATH_TO_DATA),
        clean_data_path: Path = Path(PATH_TO_CLEANED_DATA),
        vectorized_data_path: Path = Path(PATH_TO_VECTORIZED_DATA),
    ) -> None:
        # Ensure the params are paths
        if not isinstance(data_path, Path):
            raise ValueError("Data path must be a Path object.")
        if not isinstance(clean_data_path, Path):
            raise ValueError("Clean data path must be a Path object.")
        if not isinstance(vectorized_data_path, Path):
            raise ValueError("Vectorized data path must be a Path object.")
        if not data_path.exists():
            raise FileNotFoundError(f"Data path {data_path} does not exist.")
        if not clean_data_path.exists():
            raise FileNotFoundError(
                f"Clean data path {clean_data_path} does not exist."
            )
        if not vectorized_data_path.exists():
            raise FileNotFoundError(
                f"Vectorized data path {vectorized_data_path} does not exist."
            )
        self.data_path = data_path
        self.clean_data_path = clean_data_path
        self.vectorized_data_path = vectorized_data_path
        self.data = []
        self.file_types = ["pdf", "docx", "txt", "pptx", "nxml"]
        self.data_dict = {}  # Dictionary of titles and extracted content
        self.vectorized_data = {}  # Dictionary of titles and vectorized content

    def load_data(self) -> None:
        print("Loading data...")
        for root, dirs, files in os.walk(self.data_path):
            print(f"Reading {root}...")
            for file in tqdm(files):
                if file.split(".")[-1] in self.file_types:
                    self.data.append(os.path.join(root, file))
                # If .gitkeep file, ignore
                elif file == ".gitkeep":
                    continue
                # Else delete the file
                else:
                    print(f"File {file} is not in the accepted file types. Deleting...")
                    os.remove(os.path.join(root, file))

    def clean_data(self) -> None:
        print("Cleaning data...")
        for file in tqdm(self.data):
            if file.split(".")[-1] in ["pdf", "PDF"]:
                title, text = self.__clean_pdf(file)
            elif file.split(".")[-1] in ["txt", "TXT"]:
                title, text = self.__clean_txt(file)
            elif file.split(".")[-1] in ["docx", "DOCX"]:
                title, text = self.__clean_docx(file)
            elif file.split(".")[-1] in ["pptx", "PPTX"]:
                title, text = self.__clean_pptx(file)
            elif file.split(".")[-1] == "nxml":
                title, text = self.__clean_nxml(file)
            else:
                print(
                    f"File {file} is not in the accepted file types. Should have been deleted... Skipping..."
                )
                continue

            if title is None or text is None:
                print(f"Error cleaning {file}. Skipping...")
                continue
            self.data_dict[title] = text

            # Sanitize the title by replacing path separators and invalid filename characters
            invalid_chars = '<>:"/\\|?*'
            safe_title = "".join(c for c in title if c not in invalid_chars)
            safe_title = safe_title.strip()  # Remove leading/trailing whitespace

            # Save the cleaned data
            with open(
                f"{self.clean_data_path / Path(safe_title + '.txt')}",
                "w",
                encoding="utf-8",
            ) as f:
                f.write(text)

    def vectorize_data(self) -> None:
        """
        Function that vectorizes the data.
        """
        if not self.data_dict:
            print(
                "Didn't clean data, assuming it's done already and saved. Loading data..."
            )
            # Get data from the cleaned data folder
            for root, dirs, files in os.walk(self.clean_data_path):
                print(f"Reading {root}...")
                for file in tqdm(files):
                    with open(os.path.join(root, file), "r", encoding="utf-8") as f:
                        # # Print file name
                        # print(f"Reading {file}...")
                        text = f.read()
                    title = Path(file).stem
                    self.data_dict[title] = text

        if not self.data_dict:
            raise ValueError("No data to vectorize.")

        print("Vectorizing data...")
        # Vectorize the data
        for title, text in tqdm(self.data_dict.items()):
            if title == ".gitkeep":
                continue
            # print(f"Vectorizing {title}...")
            self.vectorized_data[title] = embed_text(text)

        # Save the vectorized data in json format
        with open(
            self.vectorized_data_path / Path("vectorized_data.json"),
            "w",
            encoding="utf-8",
        ) as f:
            json.dump(self.vectorized_data, f)

        return self.vectorized_data

    def load_vectorized_data(self) -> dict:
        """
        Function that loads the vectorized data.
        """
        with open(
            self.vectorized_data_path / Path("vectorized_data.json"),
            "r",
            encoding="utf-8",
        ) as f:
            self.vectorized_data = json.load(f)

        return self.vectorized_data

    def __clean_pdf(self, file: str) -> Tuple[str, str]:
        """
        Cleans a PDF file.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        print(f"Cleaning {file}...")

        # Ensure correct file type
        if file.split(".")[-1].lower() != "pdf":
            raise ValueError(f"File {file} is not a PDF file. Cannot clean as PDF.")

        # Open PDF file
        pdf = fitz.open(file)
        text = ""
        for page in pdf:
            text += page.get_text("text") + "\n\n"  # Keep paragraph breaks

        pdf.close()

        # Fix hyphenated words (e.g., "micro-\nscope" -> "microscope")
        text = re.sub(r"(\w+)-\n(\w+)", r"\1\2", text)

        # Remove unnecessary line breaks within paragraphs
        text = re.sub(r"(?<!\n)\n(?!\n)", " ", text)

        title = Path(file).stem
        return title, text

    def __clean_txt(self, file: str) -> Tuple[str, str]:
        """
        Cleans a TXT file.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] not in ["txt", "TXT"]:
            raise ValueError(f"File {file} is not a TXT file. Cannot clean as TXT.")
        # Open the TXT file
        with open(file, "r") as f:
            text = f.read()
        title = Path(file).stem
        return title, text

    def __clean_docx(self, file: str) -> Tuple[str, str]:
        """
        Cleans a DOCX file.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] not in ["docx", "DOCX"]:
            raise ValueError(f"File {file} is not a DOCX file. Cannot clean as DOCX.")
        # Open the DOCX file
        text = docx2txt.process(file)
        title = file.split("/")[-1].split(".")[0]
        return title, text

    def __clean_pptx(self, file: str) -> Tuple[str, str]:
        """
        Cleans a PPTX file.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] not in ["pptx", "PPTX"]:
            raise ValueError(f"File {file} is not a PPTX file. Cannot clean as PPTX.")
        # Open the PPTX file
        ppt = Presentation(file)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        title = Path(file).stem
        return title, text

    def __clean_nxml(self, file):
        """
        Cleans an nxml file. Specifically built to work with statpearls nxml files.

        Parameters:
        - file: str, path to the file

        Returns:
        - title: str, title of the file
        - text: str, extracted text from the file
        """
        # print(f"Cleaning {file}...")
        # Ensure right file type
        if file.split(".")[-1] != "nxml":
            raise ValueError(f"File {file} is not an nxml file. Cannot clean as nxml.")

        # Set up file extraction
        tree = et.parse(file)
        root = tree.getroot()

        # Extract the title using title-group tag
        try:
            title = root.find(".//title-group/title").text
        except AttributeError:
            print(f"Error extracting title from {file}. Skipping...")
            return None, None

        # print(f"Topic {title}...")

        # Extract tags with the sec-type as long as it does not have the value of "Continuing Education Activity"
        text = ""
        for sec in root.findall(".//sec"):
            sec_type = sec.get("sec-type")  # Access sec-type as an attribute of <sec>

            if sec_type is not None and sec_type != "Continuing Education Activity":
                for element in sec:  # Iterate over child elements in order
                    if element.tag == "title" and element.text:
                        text += element.text + "\n"
                    elif element.tag == "p" and element.text:
                        text += element.text + "\n"

        if text == "":
            print(f"Error extracting text from {file}. Skipping...")
            return None, None

        return f"StatPearls Chapter: {title}", text
